from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

# Create presentation
prs = Presentation()
prs.slide_width = Inches(10)
prs.slide_height = Inches(7.5)

def add_title_slide(title, subtitle):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    title_box = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(8), Inches(1))
    title_frame = title_box.text_frame
    title_frame.text = title
    title_p = title_frame.paragraphs[0]
    title_p.font.size = Pt(54)
    title_p.font.bold = True
    title_p.font.color.rgb = RGBColor(75, 107, 251)
    title_p.alignment = PP_ALIGN.CENTER
    
    subtitle_box = slide.shapes.add_textbox(Inches(1), Inches(4), Inches(8), Inches(0.8))
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.text = subtitle
    subtitle_p = subtitle_frame.paragraphs[0]
    subtitle_p.font.size = Pt(24)
    subtitle_p.alignment = PP_ALIGN.CENTER
    
    return slide

def add_content_slide(title, content_list):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.8))
    title_frame = title_box.text_frame
    title_frame.text = title
    title_p = title_frame.paragraphs[0]
    title_p.font.size = Pt(36)
    title_p.font.bold = True
    title_p.font.color.rgb = RGBColor(75, 107, 251)
    
    content_box = slide.shapes.add_textbox(Inches(1), Inches(1.8), Inches(8), Inches(5))
    text_frame = content_box.text_frame
    text_frame.word_wrap = True
    
    for item in content_list:
        p = text_frame.add_paragraph()
        p.text = item
        p.font.size = Pt(18)
        p.space_before = Pt(12)
        p.level = 0
    
    return slide

# Slide 1: Title
add_title_slide("EcoSystem", "Cross-Device File Sharing Platform")

# Slide 2: Project Overview
add_content_slide("Project Overview", [
    "🌐 EcoSystem - Seamless cross-device file sharing solution",
    "📱 Share files between devices using simple 6-character codes",
    "☁️ Cloud-based storage powered by Cloudinary",
    "🔒 Secure authentication with JWT tokens",
    "📊 Real-time storage tracking and management"
])

# Slide 3: Key Features
add_content_slide("Key Features", [
    "📤 File Upload - Support for images, videos, documents (up to 500MB)",
    "🔑 Access Codes - Generate unique 6-character codes for sharing",
    "📱 QR Code Generation - Scan to retrieve files instantly",
    "💬 Text Sharing - Share text content alongside files",
    "⬇️ Cross-Device Retrieval - Access files from any device",
    "📊 Storage Dashboard - Track usage with 25GB free storage"
])

# Slide 4: Technology Stack
add_content_slide("Technology Stack", [
    "Frontend:",
    "  • React.js - Modern UI framework",
    "  • Framer Motion - Smooth animations",
    "  • Axios - HTTP client",
    "  • QRCode.js - QR code generation",
    "",
    "Backend:",
    "  • Node.js + Express.js - REST API",
    "  • JWT - Authentication",
    "  • Bcrypt - Password hashing",
    "  • Cloudinary - File storage & CDN"
])

# Slide 5: Architecture
add_content_slide("System Architecture", [
    "Frontend (React) → Vercel Hosting",
    "  ↓",
    "Backend API (Express) → Render Hosting",
    "  ↓",
    "Cloudinary Storage → File CDN",
    "",
    "Authentication Flow:",
    "  • User signup/login with JWT tokens",
    "  • Persistent user storage in JSON file",
    "  • Secure password hashing with bcrypt"
])

# Slide 6: How It Works
add_content_slide("How It Works", [
    "1️⃣ Upload Process:",
    "   • User selects files or enters text",
    "   • System generates unique 6-character code",
    "   • Files uploaded to Cloudinary with code tag",
    "   • QR code generated for easy sharing",
    "",
    "2️⃣ Retrieval Process:",
    "   • User enters access code or scans QR",
    "   • System searches Cloudinary by code tag",
    "   • Files displayed with download options"
])

# Slide 7: User Interface
add_content_slide("User Interface Highlights", [
    "🏠 Home Page - Welcome screen with project intro",
    "📊 Dashboard - File management and storage tracking",
    "📤 Upload Modal - Drag-and-drop file selection",
    "🔍 Retrieve Modal - Code entry with validation",
    "📱 Mobile Responsive - Optimized for all devices",
    "🎨 Modern Design - Clean, intuitive interface"
])

# Slide 8: Security Features
add_content_slide("Security & Privacy", [
    "🔐 JWT Authentication - Secure token-based auth",
    "🔒 Password Hashing - Bcrypt encryption",
    "🌐 CORS Protection - Controlled API access",
    "👤 User Isolation - Personal file management",
    "🔑 Unique Access Codes - Random 6-char generation",
    "⏰ Session Management - 24-hour token expiry"
])

# Slide 9: Deployment
add_content_slide("Deployment Details", [
    "Production URLs:",
    "  • Frontend: https://eco-stemm-ewx3.vercel.app",
    "  • Backend: https://ecosystem-file-share-1.onrender.com",
    "",
    "Hosting Platforms:",
    "  • Vercel - Frontend deployment with auto-deploy",
    "  • Render - Backend API with persistent storage",
    "  • Cloudinary - 25GB free cloud storage",
    "",
    "Environment:",
    "  • Node.js runtime",
    "  • Production-ready configuration"
])

# Slide 10: Storage Management
add_content_slide("Storage Management", [
    "📊 Real-time Usage Tracking",
    "  • Visual progress bar",
    "  • Used vs. Available display",
    "  • 25GB free storage limit",
    "",
    "🗑️ Developer Tools",
    "  • Clear all storage (dev@secureshare.com only)",
    "  • Storage analytics",
    "  • File management dashboard"
])

# Slide 11: User Accounts
add_content_slide("User Accounts", [
    "Test Account:",
    "  • Email: test@test.com",
    "  • Password: test123",
    "",
    "Developer Account:",
    "  • Email: dev@secureshare.com",
    "  • Password: dev123",
    "  • Additional privileges: Storage management",
    "",
    "Features:",
    "  • Persistent user storage",
    "  • Secure authentication",
    "  • Personal file tracking"
])

# Slide 12: Recent Improvements
add_content_slide("Recent Optimizations", [
    "✅ Fixed mobile keyboard input lag",
    "✅ Optimized upload speed (500MB support)",
    "✅ Improved download functionality",
    "✅ Enhanced error handling",
    "✅ Added upload progress tracking",
    "✅ Mobile-responsive design",
    "✅ QR code persistence management",
    "✅ Cross-device file retrieval"
])

# Slide 13: Future Enhancements
add_content_slide("Future Roadmap", [
    "🚀 Planned Features:",
    "  • File expiration settings",
    "  • Password-protected shares",
    "  • Bulk file operations",
    "  • File preview functionality",
    "  • Share history tracking",
    "  • Email notifications",
    "  • Advanced search filters",
    "  • Multi-language support"
])

# Slide 14: Technical Challenges Solved
add_content_slide("Challenges Overcome", [
    "⚡ Input Lag Issue:",
    "  • Problem: Mobile keyboard vanishing on input",
    "  • Solution: React.memo + ref-based state management",
    "",
    "📤 Upload Performance:",
    "  • Problem: Slow uploads for large files",
    "  • Solution: Chunked uploads + timeout handling",
    "",
    "⬇️ Download Errors:",
    "  • Problem: CORS issues with blob downloads",
    "  • Solution: Direct link downloads"
])

# Slide 15: Thank You
slide = prs.slides.add_slide(prs.slide_layouts[6])
title_box = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(8), Inches(1.5))
title_frame = title_box.text_frame
title_frame.text = "Thank You!"
title_p = title_frame.paragraphs[0]
title_p.font.size = Pt(54)
title_p.font.bold = True
title_p.font.color.rgb = RGBColor(75, 107, 251)
title_p.alignment = PP_ALIGN.CENTER

subtitle_box = slide.shapes.add_textbox(Inches(1), Inches(4.5), Inches(8), Inches(1))
subtitle_frame = subtitle_box.text_frame
subtitle_frame.text = "EcoSystem - Share Files Seamlessly"
subtitle_p = subtitle_frame.paragraphs[0]
subtitle_p.font.size = Pt(24)
subtitle_p.alignment = PP_ALIGN.CENTER

contact_box = slide.shapes.add_textbox(Inches(1), Inches(5.5), Inches(8), Inches(0.8))
contact_frame = contact_box.text_frame
contact_frame.text = "https://eco-stemm-ewx3.vercel.app"
contact_p = contact_frame.paragraphs[0]
contact_p.font.size = Pt(18)
contact_p.font.color.rgb = RGBColor(100, 100, 100)
contact_p.alignment = PP_ALIGN.CENTER

# Save presentation
prs.save('EcoSystem_Presentation.pptx')
print("✅ PowerPoint presentation created: EcoSystem_Presentation.pptx")
